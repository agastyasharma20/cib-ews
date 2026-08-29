"""
Builds docs/CIB_EWS_Presentation.pptx — an interview-ready walkthrough deck
for the CIB Early Warning System project.

Run with:
    python docs/presentation/build_presentation.py

Design: 16:9 widescreen, a small consistent template (title bar, accent
rule, page number) applied via helper functions rather than repeated
boilerplate per slide. Palette is navy + red, in the spirit of HDFC Bank's
brand colors (a visual nod to the target institution) — deliberately NOT
their actual logo or an exact brand-guideline match, since this deck is
shared publicly and colors, unlike a logo mark, aren't a trademarked asset
to reproduce without permission. Speaker notes are attached to the slides
that need real talking points, not just bullet text, so this deck is
usable to actually present from, not just read.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent.parent
FIGURES = ROOT / "docs" / "figures"
RESULTS_FIGURES = ROOT / "results" / "figures"
OUTPUT_FILE = ROOT / "docs" / "CIB_EWS_Presentation.pptx"

# --- Palette (matches app/dashboard.py's HDFC-inspired navy/red theme) ---
NAVY = RGBColor(0x0B, 0x24, 0x47)
NAVY_LIGHT = RGBColor(0x19, 0x37, 0x6D)
RED = RGBColor(0xE4, 0x00, 0x2B)
RED_DARK = RGBColor(0xB3, 0x00, 0x1F)
RED_PALE = RGBColor(0xFF, 0xC1, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0B, 0x24, 0x47)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_TEXT = RGBColor(0x92, 0x40, 0x0E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

AUTHOR_NAME = "Agastya Sharma"
AUTHOR_EMAIL = "work.agastya20@gmail.com"
REPO_URL = "github.com/agastyasharma20/cib-ews"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def _set_font(run, size=18, color=INK, bold=False, italic=False, font_name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name


def add_slide() -> "Slide":
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text: str) -> None:
    """Attaches real speaker notes — talking points to actually present
    from, not a restatement of the slide's bullets."""
    slide.notes_slide.notes_text_frame.text = text


def add_background(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_titlebar(slide, title: str, kicker: str | None = None):
    """Standard content-slide header: navy title bar + red accent rule."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.15), SLIDE_W, Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()
    accent.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    if kicker:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = kicker.upper()
        _set_font(r0, size=12, color=RED_PALE, bold=True)
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = title
        _set_font(r1, size=26, color=WHITE, bold=True)
    else:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_font(r, size=28, color=WHITE, bold=True)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide, page_num: int):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(9), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"CIB Early Warning System  •  {AUTHOR_NAME}  •  {REPO_URL}"
    _set_font(r, size=9, color=MUTED)

    tb2 = slide.shapes.add_textbox(Inches(12.3), Inches(7.15), Inches(0.6), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_num)
    _set_font(r2, size=9, color=MUTED)


def add_bullets(slide, items, left=0.7, top=1.55, width=11.9, height=5.3, size=18, gap_before=10):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(gap_before)
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = f"{head}  "
            _set_font(r1, size=size, color=RED_DARK, bold=True)
            r2 = p.add_run()
            r2.text = body
            _set_font(r2, size=size, color=INK)
        else:
            bullet = "•  " if not str(item).startswith("  ") else "‒  "
            r = p.add_run()
            r.text = bullet + str(item).strip()
            _set_font(r, size=size, color=INK)
    return tb


def add_table(slide, headers, rows, left=0.7, top=1.7, width=11.9, height=None, col_widths=None, font_size=13):
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = height or Inches(0.5 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_font(p.runs[0], size=font_size, color=WHITE, bold=True)

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if r % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            _set_font(p.runs[0], size=font_size, color=INK, bold=(c == 0))
    return shape


def add_image_centered(slide, image_path: Path, top=1.5, max_width=11.9, max_height=5.4):
    from PIL import Image

    with Image.open(image_path) as img:
        w, h = img.size
    aspect = w / h
    width = max_width
    height = width / aspect
    if height > max_height:
        height = max_height
        width = height * aspect
    left = (13.333 - width) / 2
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_caption(slide, text, top=7.0):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_font(r, size=11, color=MUTED, italic=True)


def add_step_boxes(slide, steps, top=1.7, box_h=2.2):
    """3 (or n) equal-width colored boxes across the slide — used for the
    WATCH/EXPLAIN/ACT flow and similarly simple visual sequences."""
    n = len(steps)
    gap = 0.35
    total_w = 11.9
    box_w = (total_w - gap * (n - 1)) / n
    colors = [NAVY, RED_DARK, NAVY_LIGHT, RED]
    for i, (title, body) in enumerate(steps):
        x = 0.7 + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(10); tf.margin_right = Pt(10)
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = title
        _set_font(r0, size=19, color=WHITE, bold=True)
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(6)
        r1 = p1.add_run(); r1.text = body
        _set_font(r1, size=13, color=RGBColor(0xE2, 0xE8, 0xF0))


def add_disclaimer_strip(slide, top=1.28):
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(top), Inches(12.5), Inches(0.4))
    strip.fill.solid()
    strip.fill.fore_color.rgb = AMBER_BG
    strip.line.color.rgb = RGBColor(0xFD, 0xE6, 0x8A)
    strip.line.width = Pt(0.75)
    tf = strip.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "⚠ Independent portfolio project on 100% SYNTHETIC data — not affiliated with or endorsed by HDFC Bank"
    _set_font(r, size=11, color=AMBER_TEXT, bold=True)


# ===========================================================================
# Slide 1 — Title
# ===========================================================================
s = add_slide()
add_background(s, NAVY)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), SLIDE_W, Pt(4))
accent.fill.solid(); accent.fill.fore_color.rgb = RED; accent.line.fill.background(); accent.shadow.inherit = False

logo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.87), Inches(1.15), Inches(1.6), Inches(1.6))
logo.fill.solid(); logo.fill.fore_color.rgb = WHITE
logo.line.color.rgb = RED; logo.line.width = Pt(2.5)
logo.shadow.inherit = False
lp = logo.text_frame.paragraphs[0]
lp.alignment = PP_ALIGN.CENTER
lr = lp.add_run(); lr.text = "🏦"
_set_font(lr, size=44, color=RED)

tb = s.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(1.3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "CIB Early Warning System"
_set_font(r, size=44, color=WHITE, bold=True)

tb2 = s.shapes.add_textbox(Inches(1), Inches(4.75), Inches(11.33), Inches(0.7))
tf2 = tb2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "A 30–90 Day Silent-Deterioration Risk Model for a Current Account / CIB Portfolio"
_set_font(r2, size=18, color=RGBColor(0xCB, 0xD5, 0xE1))

tb3 = s.shapes.add_textbox(Inches(1), Inches(5.9), Inches(11.33), Inches(0.5))
p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = f"{AUTHOR_NAME}  •  Data Science Internship Portfolio Project  •  {REPO_URL}"
_set_font(r3, size=13, color=RED_PALE)

tb4 = s.shapes.add_textbox(Inches(1), Inches(6.6), Inches(11.33), Inches(0.5))
p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
r4 = p4.add_run(); r4.text = "⚠ Built entirely on synthetic data — not affiliated with or endorsed by HDFC Bank"
_set_font(r4, size=12, color=RGBColor(0xFC, 0xD3, 0x4D), italic=True)

set_notes(s, (
    "Open by naming the real problem in one sentence before touching the deck: 'HDFC's CA customers don't "
    "close accounts when they're unhappy — they quietly move activity to a competitor while the account stays "
    "open.' Then say this is a from-scratch, end-to-end project built to demonstrate how you'd approach that "
    "as a data scientist, on synthetic data since real data isn't available outside the bank. Mention the "
    "navy/red theme is a deliberate nod to HDFC's brand, not their actual logo — shows you thought about "
    "professionalism and IP boundaries, which is itself a good signal to a bank's reviewers."
))


# ===========================================================================
# Slide 2 — Agenda
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Agenda")
agenda = [
    "1. Business problem — why silent deterioration, not closure",
    "2. In plain English — the 60-second version",
    "3. Approach — end-to-end pipeline, 9 phases",
    "4. Synthetic data & the deterioration index",
    "5. Feature engineering — 5 signal groups",
    "6. Model results — baseline → tuned core model",
    "7. Explainability — SHAP reason codes",
    "8. Risk sizing & actions — PFaR + RM recommendations",
    "9. Survival analysis — estimating WHEN, not just IF",
    "10. Graph feature experiment — an honest negative result",
    "11. The dashboard",
    "12. Why this matters for HDFC, and what's next",
]
add_bullets(s, agenda, size=17, gap_before=10)
add_footer(s, 2)


# ===========================================================================
# Slide 3 — Business Problem
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "The Business Problem", kicker="Why this project exists")
add_bullets(s, [
    ("Not closure —", "large CA/CIB customers keep the account technically open while quietly shifting balances, transaction flow, payroll, and trade activity to a competing bank."),
    ("Reactive today —", "RMs typically act only once balance erosion is already visible in monthly reports — by then, most of the wallet share is usually already gone."),
    ("The ask —", "a system that flags this kind of silent deterioration 30–90 days in advance, explains WHY, estimates roughly WHEN, sizes the exposure, and recommends a specific action."),
], size=19, gap_before=18)
add_footer(s, 3)
set_notes(s, (
    "This is the slide to slow down on. Emphasize 'technically open' — that's the whole insight the project "
    "is built around. If asked 'why not just watch balances', that's your cue for the next slide's framing: "
    "balance is the LAST thing to move, and the whole value of an early-warning system is watching what "
    "moves BEFORE it."
))


# ===========================================================================
# Slide 4 — In Plain English
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "In Plain English", kicker="Same idea, no jargon — the 60-second version")
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.9), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("Imagine a company banks with HDFC. They don't close the account — they just quietly start "
          "routing payments, salaries, and trade deals through a rival bank instead. By the time the "
          "balance visibly drops, they're already mostly gone.")
_set_font(r, size=17, color=INK, italic=True)
add_step_boxes(s, [
    ("1. WATCH", "Every month, compare each customer's recent behavior to their own history and to similar peers."),
    ("2. EXPLAIN", "If something looks like quiet drifting-away, say exactly which behaviors are driving that — in plain English."),
    ("3. ACT", "Estimate roughly when it'll get serious, size the ₹ at risk, and hand the RM one specific action."),
], top=2.75, box_h=2.3)
add_caption(s, "The rest of this deck is how each of those three steps was actually built and validated.", top=5.5)
add_footer(s, 4)
set_notes(s, (
    "Use this slide if the audience includes anyone non-technical (a hiring manager, an HR panelist). It's "
    "your elevator pitch — you should be able to say this paragraph from memory without looking at the "
    "slide. Everything after this slide is proving each of the three boxes actually works."
))


# ===========================================================================
# Slide 5 — Approach / Pipeline
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "End-to-End Approach", kicker="9 phases, one pipeline")
phases = [
    "Synthetic Data", "Labeling & Features", "Baseline Model",
    "Tuned Core Model", "Explainability (SHAP)", "PFaR + RM Actions",
    "Survival Analysis", "Graph Feature", "Dashboard",
]
cols = 3
box_w, box_h = 3.7, 1.15
gap_x, gap_y = 0.3, 0.35
start_x, start_y = 0.85, 1.7
for i, name in enumerate(phases):
    row, col = divmod(i, cols)
    x = start_x + col * (box_w + gap_x)
    y = start_y + row * (box_h + gap_y)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY if i % 2 == 0 else RED_DARK
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = f"{i + 1}. "
    _set_font(r1, size=15, color=WHITE, bold=True)
    r2 = p.add_run(); r2.text = name
    _set_font(r2, size=15, color=WHITE, bold=True)
add_caption(s, "Every stage reads/writes through src/config.py's file paths — a real data source swaps in at ONE seam, nothing downstream rebuilds.")
add_footer(s, 5)


# ===========================================================================
# Slide 6 — Synthetic Data & Cohorts
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Synthetic Data: 4 Planted Cohorts", kicker="Phase 1 — with hidden ground truth for validation")
add_image_centered(s, FIGURES / "amb_trajectory_by_cohort.png", top=1.4, max_height=5.0)
add_caption(s, "5,000 customers × 18 months. stable (55%) / gradual (20%) / sudden (10%) / seasonal_false_positive (15%, deliberately NOT deterioration).", top=6.55)
add_footer(s, 6)


# ===========================================================================
# Slide 7 — Deterioration Index & Label Validation
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "The Deterioration Index", kicker="Phase 2 — percentile-ranked, not a flat threshold")
add_bullets(s, [
    "Composite score: AMB decline (35%) + transaction decline (30%) + digital decline (15%) + payroll/trade decline (20%)",
    "Each component ranked against SAME-SEGMENT, SAME-MONTH peers — cancels out size differences and shared seasonality",
    "A breach reverting within 2 months AND uncorroborated by other signals is filtered out as a harmless dip",
], size=17, top=1.5, height=2.3, gap_before=10)

add_table(
    s,
    headers=["True cohort (hidden ground truth)", "Flag rate"],
    rows=[
        ["sudden_deterioration", "100.0%"],
        ["gradual_deterioration", "97.5%"],
        ["seasonal_false_positive", "1.7%"],
        ["stable", "0.1%"],
    ],
    left=0.85, top=4.0, width=6.0, col_widths=[4.3, 1.7],
)
add_bullets(s, [
    "Precision 98.9%  •  Recall 98.3%",
    "Median detection lag vs. true onset: 2 months (sudden), 4 months (gradual)",
], left=7.3, top=4.15, width=5.3, height=2.6, size=16, gap_before=14)
add_footer(s, 7)
set_notes(s, (
    "If asked 'why not just threshold the balance': explain that an SME's balance naturally swings more "
    "than a Large Corporate's, and quarter-end window dressing moves everyone's balance at once — ranking "
    "against same-segment, same-month peers cancels both effects out. The seasonal_false_positive row (1.7% "
    "flagged) is the proof that a naive dip-detector would have failed here."
))


# ===========================================================================
# Slide 8 — Feature Groups
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Feature Engineering — 5 Groups, 22 Features", kicker="Phase 2b — all leakage-safe (trailing data only)")
add_table(
    s,
    headers=["Group", "# Features", "Example"],
    rows=[
        ["1. Balance & Liquidity", "6", "amb_trend_slope_3m, seasonal-adjusted deviation"],
        ["2. Transaction & Digital Activity", "5", "txn value trend, digital channel share"],
        ["3. Product & Wallet-Share", "4", "payroll regularity score, trade utilization trend"],
        ["4. Network & Counterparty", "3", "competitor-bank transaction share, HHI"],
        ["5. Relationship & Engagement", "4", "tenure, complaint/service-ticket trend"],
    ],
    left=0.7, top=1.6, width=11.9, col_widths=[4.6, 1.5, 5.8], font_size=15,
)
add_caption(s, "Same grouping reused end-to-end: reason codes (Phase 5) and the RM-action mapping (Phase 6) key off these exact 5 groups.", top=5.1)
add_footer(s, 8)


# ===========================================================================
# Slide 9 — Model Results
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Model Results — Baseline → Tuned Core Model", kicker="Phases 3–5 — 90-day horizon shown")
add_table(
    s,
    headers=["Model", "ROC-AUC", "PR-AUC", "Top-decile lift"],
    rows=[
        ["Logistic Regression (baseline)", "0.935", "0.841", "5.85x"],
        ["LightGBM (tuned)", "0.951", "0.854", "6.16x"],
        ["XGBoost (tuned) — CHOSEN", "0.958", "0.903", "6.29x"],
    ],
    left=0.7, top=1.55, width=11.9, col_widths=[5.5, 2.1, 2.1, 2.2], font_size=16,
)
add_bullets(s, [
    "Chosen on PR-AUC, not ROC-AUC — at a ~6-15% positive rate, ROC-AUC can look strong while precision in the actionable top slice is mediocre.",
    "LightGBM's PR-AUC looked competitive, but its precision/recall COLLAPSED TO 0 at the default 0.5 threshold — miscalibrated probabilities, caught by checking beyond the headline metric.",
], left=0.7, top=3.35, width=11.9, height=1.6, size=16, gap_before=12)
add_image_centered(s, RESULTS_FIGURES / "core_gbm_xgboost_deteriorates_in_90d_roc_pr_lift.png", top=5.05, max_width=11.9, max_height=2.0)
add_footer(s, 9)
set_notes(s, (
    "The LightGBM 0.000/0.000 precision/recall row is a great interview story: it shows you don't just read "
    "the headline metric and stop — you checked a second angle (a fixed threshold) and caught something the "
    "PR-AUC alone hid. Lead with that if asked 'tell me about a time you caught your own mistake / a subtle bug'."
))


# ===========================================================================
# Slide 10 — Explainability
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Explainability — SHAP Reason Codes", kicker="Phase 5 — from a score to a story")
add_bullets(s, [
    "SHAP TreeExplainer computes PER-CUSTOMER attributions — not just global feature importance",
    "Raw features grouped into plain-language driver labels BEFORE ranking (3 distinct stories, not 1 story 3 times)",
    "Ranked by signed SHAP sum, not magnitude — a protective factor never crowds out the actual risk drivers",
], size=17, top=1.5, height=2.2, gap_before=10)
add_table(
    s,
    headers=["Risk", "90d score", "Top 3 reason codes"],
    rows=[
        ["High", "1.00", "Falling balances → Lower digital activity → Reduced payroll activity"],
        ["Medium", "0.50", "Declining credits → Concentrating counterparties → Declining transactions"],
        ["Low", "0.00", "Narrowing product relationship → Rising complaints → Concentrating counterparties"],
    ],
    left=0.7, top=3.9, width=11.9, col_widths=[1.3, 1.3, 9.3], font_size=14,
)
add_footer(s, 10)


# ===========================================================================
# Slide 11 — PFaR & RM Actions
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Risk Sizing & Recommended Actions", kicker="Phase 6 — PFaR (Probability-weighted Funds at Risk)")
tb = s.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.6), Inches(0.7))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "PFaR  =  90-day risk score  ×  (current balance × empirically-calibrated expected decline %)"
_set_font(r, size=20, color=RED_DARK, bold=True)
add_bullets(s, [
    "Decline % calibrated from OUR OWN historically-confirmed deterioration cases — never from the hidden ground truth (that would be leakage a real deployment could never replicate)",
    "Decomposed into a driver TYPE — liquidity / relationship / competitor — from the customer's top reason code",
    "RM action mapping: a swappable rules-based lookup (top reason code → action), interface shaped for a future contextual-bandit upgrade once real RM-outcome data exists",
], top=2.5, height=2.6, size=17, gap_before=14)
add_table(
    s,
    headers=["Reason code", "Recommended action"],
    rows=[
        ["Falling balances", "Discuss liquidity requirements"],
        ["Rising competitor-transfer share", "Proactive wallet-share conversation / pricing review"],
        ["Reduced payroll activity", "Assess payroll migration risk"],
    ],
    left=0.85, top=5.05, width=11.6, height=Inches(1.8), col_widths=[5.3, 6.3], font_size=14,
)
add_footer(s, 11)


# ===========================================================================
# Slide 12 — Survival Analysis
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Survival Analysis — Estimating WHEN", kicker="Phase 7 — time-to-deterioration")
add_bullets(s, [
    ("Tried first:", "CoxPHFitter — the standard approach."),
    ("Checked, not assumed:", "check_assumptions() flagged 11 of 24 covariates (46%) violating proportional hazards — pervasive, not borderline."),
    ("Switched to:", "Random Survival Forest — no such assumption. Concordance index 0.659 vs. Cox's 0.591."),
], left=0.7, top=1.55, width=6.1, height=4.0, size=16, gap_before=14)
add_image_centered(s, RESULTS_FIGURES / "survival_example_curves.png", top=1.5, max_width=5.9, max_height=4.6)
add_caption(s, "0.659 is modest next to the classifiers' 0.90+ AUC — expected: this model gets ONE early snapshot per customer, not fresh features every month.", top=6.35)
add_footer(s, 12)
set_notes(s, (
    "This is your best 'I check my assumptions instead of trusting a library's default' story. Cox is the "
    "textbook answer; running check_assumptions() and actually reading the output (11/24 covariates failing) "
    "before trusting the hazard ratios is the part worth narrating slowly if asked about statistical rigor."
))


# ===========================================================================
# Slide 13 — Graph Feature: Honest Negative Result
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Graph Feature — An Honest Negative Result", kicker="Phase 8 (stretch) — networkx wallet-leakage graph")
add_table(
    s,
    headers=["Horizon", "Variant", "ROC-AUC", "PR-AUC"],
    rows=[
        ["90d", "Before (no graph)", "0.9585", "0.9027"],
        ["90d", "After (+ graph feature)", "0.9580", "0.9019"],
    ],
    left=0.7, top=1.6, width=8.0, col_widths=[1.6, 3.4, 1.5, 1.5], font_size=15,
)
add_bullets(s, [
    ("No lift:", "flat to very slightly worse at every horizon; ranks 14th–25th of 27 features by SHAP."),
    ("Root cause, confirmed:", "the graph feature correlates 0.88–0.91 with an existing Phase 2 feature built from the same transactions — methodologically distinct, mostly redundant information."),
    ("Why keep it:", "shows the negative result was investigated and explained, not hidden — and the graph is the right extension point once real counterparty-entity data exists."),
], top=3.3, height=3.4, size=16, gap_before=14)
add_footer(s, 13)
set_notes(s, (
    "Don't skip or rush this slide — a negative result you investigated and explained is a STRONGER signal "
    "of maturity than another slide of things that worked. If a reviewer asks 'what didn't work', this is "
    "the answer, and you already have the root cause (0.88-0.91 correlation) ready, not just 'it didn't help'."
))


# ===========================================================================
# Slide 14 — Dashboard
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "The Dashboard", kicker="Phase 9 — RM Cockpit + Portfolio View")
add_image_centered(s, FIGURES / "dashboard_rm_cockpit.png", top=1.4, max_height=5.5)
add_caption(s, "RM Cockpit: PFaR-ranked, filterable, with plain-language reason codes and an inline survival-curve sparkline per customer.", top=7.0)
add_footer(s, 14)


# ===========================================================================
# Slide 15 — Why This Could Matter for HDFC
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Why This Could Matter for HDFC", kicker="Business framing, not a real revenue estimate")
add_bullets(s, [
    ("CIB relationships are high-value —", "large balances, deep cross-sell, and expensive to win back once fully gone. Even a modest improvement in EARLY detection compounds at portfolio scale."),
    ("Shifts RM effort from reactive to prioritized —", "instead of a monthly balance report an RM notices too late, a ranked worklist with a reason and a specific action."),
    ("The ₹8,570 Cr flagged exposure on this synthetic 5,000-customer book isn't a real estimate —", " it's a demonstration of the SHAPE of the calculation a real portfolio-scale PFaR number would take."),
    ("The real ask:", "a chance to do this analysis on real (anonymized/sampled) CIB data, where every threshold and calibration in this deck would be re-tuned, not reused blindly (see the Path to Production section of docs/methodology.md)."),
], size=16, gap_before=14, height=5.3)
add_footer(s, 15)
set_notes(s, (
    "Be careful with this slide — do not let it sound like you're claiming to know HDFC's real numbers. The "
    "point is showing you think about business impact, not just model metrics, while being explicit that the "
    "rupee figure is illustrative. If pressed, say plainly: 'I don't have HDFC's real numbers — this shows "
    "the mechanism, and I'd want the chance to run it on real data.'"
))


# ===========================================================================
# Slide 16 — Innovation & What's Next
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Innovation & What's Next", kicker="How this grows from a portfolio project into a stronger one")
add_bullets(s, [
    ("Close the one gap myself:", "an automated test suite (label thresholds, feature leakage-safety, PFaR logic) — the one thing this project doesn't have yet."),
    ("Contextual bandit RM-action recommender:", "the action-engine interface is already shaped for it — needs real RM-outcome data to learn from, not more code."),
    ("The REAL wallet-leakage graph:", "customer↔customer, via shared vendors/directors, once entity-level counterparty data exists — Phase 8 showed exactly why the bank-level version fell short."),
    ("Model monitoring:", "track PR-AUC and calibration drift over time as real data accumulates, not just a one-time evaluation."),
    ("NLP on RM call notes / service tickets:", "a 6th feature group, using text the current data doesn't include."),
], size=15.5, gap_before=12, height=5.3)
add_footer(s, 16)
set_notes(s, (
    "This slide is deliberately about what you'd do NEXT, not more of what's already done — use it to show "
    "self-awareness of the project's edges and genuine ideas for extending it, which is exactly what an "
    "internship interviewer wants to see: can this person keep improving something after the first version "
    "ships, not just finish an assignment."
))


# ===========================================================================
# Slide 17 — Limitations & Path to Production
# ===========================================================================
s = add_slide()
add_background(s)
add_titlebar(s, "Limitations & Path to Production", kicker="Being direct about what changes with real data")
add_bullets(s, [
    ("Known limitation:", "0.90+ AUCs are unusually strong because engineered features are correlated with how the synthetic labels were built — real data will likely perform lower."),
    ("One seam to swap:", "src/config.py's file paths are the ONLY place a real core-banking source plugs in — nothing downstream (features, models, dashboard) needs to change."),
    ("Must be recalibrated:", "DI threshold, PFaR decline %, model hyperparameters, the Cox-vs-RSF decision — all tuned on synthetic data, not assumed to transfer."),
], size=16, gap_before=14, height=5.3)
add_footer(s, 17)


# ===========================================================================
# Slide 18 — Thank You / Contact
# ===========================================================================
s = add_slide()
add_background(s, NAVY)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35), SLIDE_W, Pt(4))
accent.fill.solid(); accent.fill.fore_color.rgb = RED; accent.line.fill.background(); accent.shadow.inherit = False

tb = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.33), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You"
_set_font(r, size=40, color=WHITE, bold=True)

tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(1.2))
tf2 = tb2.text_frame
p1 = tf2.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
r1 = p1.add_run(); r1.text = AUTHOR_NAME
_set_font(r1, size=20, color=WHITE, bold=True)
p2 = tf2.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = f"{AUTHOR_EMAIL}  •  {REPO_URL}"
_set_font(r2, size=15, color=RED_PALE)

tb3 = s.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11.33), Inches(0.6))
p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Full methodology, every result, and the code: see docs/methodology.md in the repository."
_set_font(r3, size=13, color=RGBColor(0xCB, 0xD5, 0xE1), italic=True)

prs.save(OUTPUT_FILE)
print(f"Saved {OUTPUT_FILE} ({len(prs.slides)} slides)")
